from pptx import Presentation


def create_ppt(slides_data, output_file):
    prs = Presentation()

    for slide_info in slides_data:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = slide_info["title"]

        content = slide.placeholders[1]
        content.text = "\n".join(slide_info["bullets"])

        speaker_notes = slide_info.get("speaker_notes", "").strip()
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    prs.save(output_file)
